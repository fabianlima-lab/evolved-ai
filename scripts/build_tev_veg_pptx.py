#!/usr/bin/env python3
"""Build the TEV & VEG Partnership deck as a native 16:9 PowerPoint that
mirrors the HTML deck slide-for-slide, in the TEV brand palette."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

CREAM   = RGBColor(0xF7, 0xF4, 0xEE)
INK     = RGBColor(0x22, 0x22, 0x20)
TEAL    = RGBColor(0x2A, 0x7C, 0x7B)
TEALDP  = RGBColor(0x1E, 0x51, 0x50)
MINT    = RGBColor(0xB8, 0xD8, 0xD0)
MUTED   = RGBColor(0x6B, 0x6B, 0x64)
LINE    = RGBColor(0xDD, 0xD8, 0xCC)
PAPER   = RGBColor(0xFF, 0xFF, 0xFF)
CALLOUT = RGBColor(0xE9, 0xF1, 0xEC)
CREAMTX = RGBColor(0xF3, 0xF1, 0xE9)
COVSUB  = RGBColor(0xBC, 0xC9, 0xC2)

SERIF = "Georgia"
SANS  = "DM Sans"
EMU_IN = 914400
SW, SH = 13.333, 7.5
ML = 0.85
CW = SW - 2 * ML

prs = Presentation()
prs.slide_width  = Emu(int(SW * EMU_IN))
prs.slide_height = Emu(int(SH * EMU_IN))
BLANK = prs.slide_layouts[6]


def bg(slide, color):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color


def textbox(slide, x, y, w, h, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    return tf


def run(p, text, size, color, font=SANS, bold=False, italic=False):
    r = p.add_run(); r.text = text
    r.font.size = Pt(size); r.font.color.rgb = color
    r.font.name = font; r.font.bold = bold; r.font.italic = italic
    return r


def eyebrow(slide, text, y=0.55, color=TEAL):
    tf = textbox(slide, ML, y, CW, 0.35)
    run(tf.paragraphs[0], " ".join(list(text.upper())), 10, color, font=SANS, bold=True)


def slide_num(slide, text):
    tf = textbox(slide, SW - 4.85, 0.4, 4.0, 0.3)
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.RIGHT
    run(p, text, 9, MUTED, font=SANS)


def heading(slide, y, text, size=34, width=None):
    tf = textbox(slide, ML, y, width or (CW * 0.86), 1.6)
    run(tf.paragraphs[0], text, size, INK, font=SERIF)


def lede(slide, y, text, width=0.88):
    tf = textbox(slide, ML, y, CW * width, 0.9)
    run(tf.paragraphs[0], text, 12, INK, font=SANS, italic=True)


def rounded(slide, x, y, w, h, fill, border=None, radius=0.04):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                 Inches(x), Inches(y), Inches(w), Inches(h))
    shp.adjustments[0] = radius
    shp.fill.solid(); shp.fill.fore_color.rgb = fill
    if border:
        shp.line.color.rgb = border; shp.line.width = Pt(0.75)
    else:
        shp.line.fill.background()
    shp.shadow.inherit = False
    return shp


def card(slide, x, y, w, h, title, body):
    shp = rounded(slide, x, y, w, h, PAPER, border=LINE)
    tf = shp.text_frame; tf.word_wrap = True
    tf.margin_left = tf.margin_right = Inches(0.15)
    tf.margin_top = tf.margin_bottom = Inches(0.12)
    if title:
        run(tf.paragraphs[0], " ".join(list(title.upper())), 9, TEAL, font=SANS, bold=True)
        tf.paragraphs[0].space_after = Pt(4)
        bp = tf.add_paragraph()
    else:
        bp = tf.paragraphs[0]
    run(bp, body, 10, INK, font=SANS)


def bullets(slide, y, items, size=11.5, gap=6, width=0.92):
    tf = textbox(slide, ML, y, CW * width, SH - y - 0.4)
    first = True
    for lead, rest in items:
        p = tf.paragraphs[0] if first else tf.add_paragraph(); first = False
        run(p, "—  ", size, TEAL, font=SANS, bold=True)
        if lead:
            run(p, lead + " ", size, INK, font=SANS, bold=True)
        run(p, rest, size, INK, font=SANS)
        p.space_after = Pt(gap)


def callout(slide, x, y, w, h, segments):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(0.05), Inches(h))
    bar.fill.solid(); bar.fill.fore_color.rgb = TEAL; bar.line.fill.background(); bar.shadow.inherit = False
    box = rounded(slide, x + 0.05, y, w - 0.05, h, CALLOUT, radius=0.02)
    tf = box.text_frame; tf.word_wrap = True
    tf.margin_left = Inches(0.2); tf.margin_right = Inches(0.18)
    tf.margin_top = tf.margin_bottom = Inches(0.1)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    for seg, b, it in segments:
        run(p, seg, 11, INK, font=SANS, bold=b, italic=it)


def banner(slide, x, y, w, h, label, text):
    box = rounded(slide, x, y, w, h, TEALDP, radius=0.02)
    tf = box.text_frame; tf.word_wrap = True
    tf.margin_left = Inches(0.2); tf.margin_right = Inches(0.18)
    tf.margin_top = tf.margin_bottom = Inches(0.09)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    run(p, " ".join(list(label.upper())) + "    ", 10, MINT, font=SANS, bold=True)
    run(p, text, 11, CREAMTX, font=SANS)


def footnote(slide, y, text):
    tf = textbox(slide, ML, y, CW, 0.9)
    run(tf.paragraphs[0], text, 7.5, MUTED, font=SANS)


def figure_card(slide, x, y, w, h, title, big, small, sub, big_color=TEAL):
    shp = rounded(slide, x, y, w, h, PAPER, border=LINE)
    tf = shp.text_frame; tf.word_wrap = True
    tf.margin_left = tf.margin_right = Inches(0.16)
    tf.margin_top = tf.margin_bottom = Inches(0.14)
    run(tf.paragraphs[0], " ".join(list(title.upper())), 9, TEAL, font=SANS, bold=True)
    tf.paragraphs[0].space_after = Pt(4)
    p2 = tf.add_paragraph()
    run(p2, big, 30, big_color, font=SERIF)
    run(p2, small, 11, MUTED, font=SANS)
    p2.space_after = Pt(4)
    run(tf.add_paragraph(), sub, 10, INK, font=SANS)


# ---- 1 · COVER ----
s = prs.slides.add_slide(BLANK); bg(s, TEALDP)
eyebrow(s, "TEV & VEG Partnership", y=1.45, color=MINT)
run(textbox(s, ML, 1.95, CW * 0.82, 2.4).paragraphs[0],
    "The Evolved Vets Membership — for every VEGgie.", 46, CREAMTX, font=SERIF)
rule = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(ML), Inches(4.35), Inches(0.75), Inches(0.03))
rule.fill.solid(); rule.fill.fore_color.rgb = MINT; rule.line.fill.background(); rule.shadow.inherit = False
run(textbox(s, ML, 4.65, CW * 0.62, 1.4).paragraphs[0],
    "A wellness & leadership program — RACE-approved CE, live coaching, community, and recovery tools built by one of your own, for every veterinary professional at VEG.",
    15, COVSUB, font=SANS)

# ---- 2 · VISION ----
s = prs.slides.add_slide(BLANK); bg(s, CREAM)
slide_num(s, "02 · The Vision"); eyebrow(s, "The vision")
heading(s, 0.95, "The next evolution of VEG's innovation: its people.", size=30)
lede(s, 1.9, "VEG has reimagined emergency medicine for pets. This partnership extends that same innovation to the wellbeing and leadership of the people delivering that care every day — and the peer-reviewed research says it works.")
bullets(s, 2.85, [
    ("Meditation works — for burnout and for leadership development.",
     "A mindfulness course for healthcare providers significantly reduced burnout on the Maslach Burnout Inventory scale and improved mental wellbeing. A major review in The Lancet confirms these programs meaningfully reduce burnout. And when leaders practice mindfulness, their teams do better."),
    ("The need reaches the whole team.",
     "In the peer-reviewed Merck Animal Health Veterinary Team study, conducted with the AVMA, serious psychological distress was twice as prevalent among nonveterinarian practice team members as among veterinarians."),
    ("The stakes are real.",
     "Veterinarians die by suicide at 2.1× (men) and 3.5× (women) the general-population rate; veterinary technicians 5.0× (men — small-sample estimate) and 2.3× (women)."),
    ("And attrition is expensive.",
     "Replacing one employee costs one-half to two times their annual salary (Gallup), with average veterinary team turnover near 23% (AAHA)."),
], size=10.5, gap=4)
callout(s, ML, 6.1, CW * 0.96, 0.72, [
    ("This program is VEG's care for its VEGgies made tangible — part of how \"VEGgies brag about their jobs and their friends are jealous,\" and it lives the ethos: ", False, False),
    ("\"We treat pets like people and people like humans.\"", False, True),
])

# ---- 3 · MEMBERSHIP ----
s = prs.slides.add_slide(BLANK); bg(s, CREAM)
slide_num(s, "03 · The Program"); eyebrow(s, "What the membership is")
heading(s, 0.95, "Tools for the hard days. Support for the long game.", size=27)
lede(s, 1.7, "VEG gains a year-round system for building confident, energized, communication-strong leaders — made by a veterinarian, for veterinary professionals.", width=0.86)
cards3 = [
    ("TEV Connect", "A private community of veterinary professionals supporting each other's growth, expansion, and leadership."),
    ("Vet Resets", "Quick meditations made for veterinary professionals — Shift Shedding, Surgery Confidence Boost, Reconnecting With Your Why."),
    ("Live coaching", "Monthly live group coaching calls supporting each member's personal and professional evolution."),
    ("Masterclass Evolutions", "Biannual group programs — Summer and New Year — working through a new Masterclass together, week by week."),
    ("RACE-approved CE", "30 hours of RACE approval available throughout the year — wellbeing work that also counts toward licensure."),
    ("Audio archive", "A growing collection of every meditation, elevation, and integration from TEV Masterclasses, on demand."),
]
g = 0.18; cw = (CW - 2 * g) / 3; ch = 1.3; y0 = 2.4
for i, (t, b) in enumerate(cards3):
    r, c = divmod(i, 3)
    card(s, ML + c * (cw + g), y0 + r * (ch + g), cw, ch, t, b)
callout(s, ML, 5.7, CW * 0.96, 0.52, [
    ("And it never stands still: VEGgies get every ", False, False),
    ("new course, meditation, and monthly coaching call", True, False),
    (" as they're released — the membership grows all year, at no added cost to VEG.", False, False),
])
banner(s, ML, 6.42, CW * 0.96, 0.55, "VEG Bonus",
       "Hospital-wide facilitations — available to individual VEG hospital locations upon request.")

# ---- 4 · HOW IT WORKS ----
s = prs.slides.add_slide(BLANK); bg(s, CREAM)
slide_num(s, "04 · The Program"); eyebrow(s, "How it works")
heading(s, 0.95, "Made for the way VEGgies actually work.", size=30)
g = 0.2; cw = (CW - g) / 2; y0 = 2.1
card(s, ML, y0, cw, 2.75, "Created by an expert who lives it",
     "Dr. Bethany Weinheimer is a practicing ER veterinarian who knows overnights, caseload stress, and compassion fatigue firsthand. A Certified 3VQ Coach since 2019, she lectures to hospitals on culture and wellbeing, and has spent 12 years studying and applying neuroplasticity to help prevent burnout in the veterinary profession. That work became the TEV Method — a process for executive coaching centered on wellbeing and leadership development.")
card(s, ML + cw + g, y0, cw, 1.3, "Already built & running",
     "The membership is live today with paying members: community, coaching calls, meditation library, masterclasses, and RACE-approved CE all operating now.")
card(s, ML + cw + g, y0 + 1.45, cw, 1.3, "One click via Okta",
     "The membership site lives as a tile in VEG's Okta dashboard. A VEGgie clicks it and lands in the member portal — no friction between a hard shift and the tools.")
card(s, ML, y0 + 2.9, cw, 1.15, "Zero IT burden on VEG",
     "Hosting, enrollment, content, and support are all handled by the program. VEG's only role is the partnership and adding the tile.")
card(s, ML + cw + g, y0 + 2.9, cw, 1.15, "A playbook VEG already knows",
     "Like the farm-training program — an independent initiative VEG invests in because of its impact on VEGgies.")

# ---- 5 · THE MODEL ----
s = prs.slides.add_slide(BLANK); bg(s, CREAM)
slide_num(s, "05 · The Model"); eyebrow(s, "The model")
heading(s, 0.95, "A playbook VEG already knows.", size=32)
lede(s, 1.9, "Like the farm-training program: an independent initiative that VEG invests in because of its impact on VEGgies.", width=0.86)
bullets(s, 2.7, [
    ("Bethany owns and operates the program", "as an independent initiative, not VEG-branded, per your requests."),
    ("An enterprise partnership.", "Unlimited access for every VEG employee through one annual organizational investment."),
    ("One membership. All VEGgies.", "Open to every VEGgie — and the investment stays fixed as VEG grows. With 28 new hospitals opening next year, that's at least another 900 VEGgies included."),
    ("Outside sales stay separate.", "The public-facing product continues independently and will never present itself as VEG-sponsored or VEG-branded."),
    ("Built around VEGgies.", "VEG shares what its people need, and Bethany keeps building the courses and meditations with their needs in mind — reviewed together on an ongoing basis."),
], size=11.5, gap=7)

# ---- 6 · THE RETURN ----
s = prs.slides.add_slide(BLANK); bg(s, CREAM)
slide_num(s, "06 · The Return"); eyebrow(s, "The return")
heading(s, 0.95, "The returns.", size=32)
lede(s, 1.85, "Retention, leadership, and culture — value measured in the terms VEG already uses.", width=0.86)
g = 0.18; cw = (CW - 2 * g) / 3; y0 = 2.45
card(s, ML, y0, cw, 2.55, "One retained veterinarian",
     "One retained veterinarian more than covers the entire annual investment. The average practice generates ~$555,000 in revenue per veterinarian per year (AVMA, 2024), and the average DVM hire takes 15 months to fill (AAHA, 2023) — an empty seat costs roughly $46,000 a month in lost production alone.")
card(s, ML + cw + g, y0, cw, 1.55, "One leader developed",
     "A leader promoted — or kept in seat — through greater emotional regulation and growth shapes the experience of every VEGgie they lead. Leadership development compounds.")
card(s, ML + 2 * (cw + g), y0, cw, 1.85, "The scale of turnover",
     "At the industry's ~23% team turnover (AAHA), 5,200 VEGgies means roughly 1,200 team transitions a year. Preventing even 1–2% of departures would return several multiples of its annual cost.")
card(s, ML + cw + g, y0 + 1.7, cw, 1.3, "Beyond retention",
     "Lower burnout and stronger day-to-day wellbeing show up in engagement, culture, and quality of care — the outcomes VEG already measures and celebrates.")
card(s, ML + 2 * (cw + g), y0 + 2.0, cw, 1.3, "VEG's strategic priorities",
     "With growing case volumes, retaining highly skilled talent is vital — and the capabilities this membership builds are exactly what keeps people at VEG.")
footnote(s, 6.6, "Sources: AVMA Economic State of the Veterinary Profession benchmarking, 2024 ($554,982 avg. revenue per FTE veterinarian). AAHA survey, 2023 (avg. DVM time-to-hire: 15 months). AAHA turnover data as cited by the AVMA (~23% avg. annual team turnover).")

# ---- 7 · INVESTMENT ----
s = prs.slides.add_slide(BLANK); bg(s, CREAM)
slide_num(s, "07 · The Investment"); eyebrow(s, "The investment")
heading(s, 0.95, "One organizational investment. Every VEGgie included.", size=27)
lede(s, 1.75, "This enterprise partnership provides unlimited access for every VEG employee through one annual organizational investment.", width=0.86)
g = 0.2; cw = (CW - 2 * g) / 3; y0 = 2.5; ch = 2.45
figure_card(s, ML, y0, cw, ch, "Annual investment", "$58,800", " /year", "Billed monthly at $4,900 — fixed regardless of enrollment", big_color=TEAL)
figure_card(s, ML + cw + g, y0, cw, ch, "Who's included", "5,200+", " VEGgies", "Every employee, unlimited access, all new content as it's released", big_color=INK)
figure_card(s, ML + 2 * (cw + g), y0, cw, ch, "As VEG grows", "$0", " more", "Scales automatically — 28 new hospitals next year (~900 more VEGgies and 28 more facilitation opportunities), no per-user licensing, no IT lift", big_color=TEAL)
callout(s, ML, 5.25, CW * 0.96, 1.05, [
    ("Structured as a fixed annual operating agreement. Includes everything: the full membership, every new course and meditation, the monthly coaching call, hospital-wide facilitations for individual locations upon request, RACE-approved CE, Okta and app access, and Bethany's ongoing program leadership.", False, False),
])

# ---- 8 · ACCESS ----
s = prs.slides.add_slide(BLANK); bg(s, CREAM)
slide_num(s, "08 · Access"); eyebrow(s, "Access")
heading(s, 0.95, "Two ways in — at work and on the go.", size=30)
lede(s, 1.95, "VEGgies can reach the membership through a tile in VEG's Okta dashboard, or from anywhere through the Evolved Vets App.", width=0.86)
g = 0.2; cw = (CW - 2 * g) / 3; y0 = 2.9
access = [
    ("Via Okta", "A tile in VEG's Okta dashboard, right next to the tools VEGgies use every shift. One click and they're in the member portal."),
    ("Via the Evolved Vets App", "Meditations, community, courses, and coaching calls on their phone — before a shift, after a hard case, wherever they are."),
    ("A partnership with VEG IT", "Bethany will partner with VEG's technology team to stand up the tile and access — IT guides the how, TEV provides everything needed."),
]
for i, (t, b) in enumerate(access):
    card(s, ML + i * (cw + g), y0, cw, 1.95, t, b)

# ---- 9 · CLOSING ----
s = prs.slides.add_slide(BLANK); bg(s, TEALDP)
eyebrow(s, "In closing", y=2.3, color=MINT)
run(textbox(s, ML, 2.85, CW * 0.9, 2.5).paragraphs[0],
    "I'd be honored to partner with VEG in building the healthiest, most resilient veterinary workforce in the profession.",
    30, CREAMTX, font=SERIF)

out = "docs/proposals/tev-veg-partnership.pptx"
prs.save(out)
print("saved", out, "slides:", len(prs.slides._sldIdLst))
